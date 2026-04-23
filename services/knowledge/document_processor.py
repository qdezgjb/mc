"""Document Processing Service for Knowledge Space.

Author: lycosa9527
Made by: MindSpring Team

Extracts text from various file types including PDF, DOCX, images with OCR.

Copyright 2024-2025 北京思源智教科技有限公司 (Beijing Siyuan Zhijiao Technology Co., Ltd.)
All Rights Reserved
Proprietary License
"""

from pathlib import Path
from typing import Optional, Dict, Any, Tuple, List
import logging
import re
import mimetypes
import zipfile
import base64

logger = logging.getLogger(__name__)

# Try to import language detection library
try:
    from langdetect import detect, LangDetectException

    LANGDETECT_AVAILABLE = True
except ImportError:
    LANGDETECT_AVAILABLE = False
    logger.warning("[DocumentProcessor] langdetect not available, language detection disabled")

# Try to import optional dependencies
try:
    import PyPDF2

    PYPDF2_AVAILABLE = True
except ImportError:
    PYPDF2_AVAILABLE = False

try:
    import pdfplumber

    PDFPLUMBER_AVAILABLE = True
except ImportError:
    PDFPLUMBER_AVAILABLE = False

try:
    from docx import Document

    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

try:
    import httpx

    HTTPX_AVAILABLE = True
except ImportError:
    HTTPX_AVAILABLE = False

try:
    from config.settings import config

    CONFIG_AVAILABLE = True
except ImportError:
    CONFIG_AVAILABLE = False

try:
    import pytesseract
    from PIL import Image

    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False

try:
    from pptx import Presentation

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from openpyxl import load_workbook

    OPENPYXL_AVAILABLE = True
except ImportError:
    OPENPYXL_AVAILABLE = False


class DocumentProcessor:
    """
    Document processing service for text extraction.

    Supports: PDF, DOCX, TXT, MD, images (with OCR), PPTX, XLSX

    Includes file content validation using magic bytes to ensure files match their claimed type.
    """

    # Magic bytes (file signatures) for content validation
    # Format: (magic_bytes, mime_type, description)
    FILE_SIGNATURES = [
        # PDF
        (b"%PDF", "application/pdf", "PDF"),
        # DOCX (ZIP-based format, starts with PK)
        (
            b"PK\x03\x04",
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "DOCX",
        ),
        # PPTX (ZIP-based format)
        (
            b"PK\x03\x04",
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "PPTX",
        ),
        # XLSX (ZIP-based format)
        (
            b"PK\x03\x04",
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "XLSX",
        ),
        # Images
        (b"\xff\xd8\xff", "image/jpeg", "JPEG"),
        (b"\x89PNG\r\n\x1a\n", "image/png", "PNG"),
        (b"GIF87a", "image/gif", "GIF87a"),
        (b"GIF89a", "image/gif", "GIF89a"),
        (b"BM", "image/bmp", "BMP"),
        (b"II*\x00", "image/tiff", "TIFF (little-endian)"),
        (b"MM\x00*", "image/tiff", "TIFF (big-endian)"),
        # Text files (no magic bytes, but we can check for UTF-8 BOM)
        (b"\xef\xbb\xbf", "text/plain", "UTF-8 BOM"),
    ]

    # Maximum magic bytes length to read
    MAX_MAGIC_LEN = max(len(sig[0]) for sig in FILE_SIGNATURES)

    def __init__(self):
        """Initialize document processor."""
        self.supported_types = {
            "application/pdf": self._extract_pdf,
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document": self._extract_docx,
            "text/plain": self._extract_text,
            "text/markdown": self._extract_text,
            "image/jpeg": self._extract_image,
            "image/png": self._extract_image,
            "image/jpg": self._extract_image,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation": self._extract_pptx,
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet": self._extract_xlsx,
        }

    def validate_file_content(self, file_path: str, expected_mime_type: str) -> Tuple[bool, Optional[str]]:
        """
        Validate file content using magic bytes (file signatures).

        Args:
            file_path: Path to file
            expected_mime_type: Expected MIME type

        Returns:
            Tuple of (is_valid, detected_mime_type)
            - is_valid: True if file content matches expected type
            - detected_mime_type: Detected MIME type from magic bytes (None if not detected)
        """
        try:
            with open(file_path, "rb") as f:
                header = f.read(self.MAX_MAGIC_LEN)

            if not header:
                return False, None

            # Check magic bytes
            detected_type = None
            for magic_bytes, mime_type, _description in self.FILE_SIGNATURES:
                if header.startswith(magic_bytes):
                    detected_type = mime_type
                    break

            # Special handling for ZIP-based formats (DOCX, PPTX, XLSX)
            # They all start with PK\x03\x04, so we need to check the internal structure
            if header.startswith(b"PK\x03\x04"):
                # Check for Office Open XML structure
                # DOCX contains word/document.xml
                # PPTX contains ppt/presentation.xml
                # XLSX contains xl/workbook.xml
                try:
                    with zipfile.ZipFile(file_path, "r") as zip_file:
                        file_list = zip_file.namelist()

                        # Check for DOCX
                        if any("word/document.xml" in f for f in file_list):
                            detected_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                        # Check for PPTX
                        elif any("ppt/presentation.xml" in f for f in file_list):
                            detected_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
                        # Check for XLSX
                        elif any("xl/workbook.xml" in f for f in file_list):
                            detected_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
                        else:
                            # It's a ZIP but not a recognized Office format
                            detected_type = "application/zip"
                except Exception as e:
                    logger.warning("[DocumentProcessor] Failed to inspect ZIP structure: %s", e)
                    # If ZIP inspection fails, assume it matches if expected type is ZIP-based
                    if expected_mime_type in [
                        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    ]:
                        detected_type = expected_mime_type

            # For text files, we can't reliably detect from magic bytes alone
            # But we can check if it's valid UTF-8
            if expected_mime_type in ["text/plain", "text/markdown"]:
                try:
                    with open(file_path, "r", encoding="utf-8") as f:
                        f.read(1024)  # Try to read first 1KB
                    detected_type = expected_mime_type  # Assume valid if readable as UTF-8
                except UnicodeDecodeError:
                    return False, None

            # Validate match
            if detected_type:
                # Allow some flexibility for image types
                if expected_mime_type in ["image/jpeg", "image/jpg"] and detected_type == "image/jpeg":
                    return True, detected_type
                if expected_mime_type == detected_type:
                    return True, detected_type
                else:
                    logger.warning(
                        "[DocumentProcessor] File content mismatch: expected %s, detected %s",
                        expected_mime_type,
                        detected_type,
                    )
                    return False, detected_type

            # If no magic bytes match but file exists, allow it (might be text or other format)
            # But log a warning
            logger.warning(
                "[DocumentProcessor] Could not detect file type from magic bytes: %s",
                file_path,
            )
            return True, None  # Allow processing but warn

        except Exception as e:
            logger.error("[DocumentProcessor] Failed to validate file content: %s", e)
            return False, None

    def extract_references(self, text: str, document_id: int) -> List[Dict[str, Any]]:
        """
        Extract references and citations from document text.

        Args:
            text: Document text
            document_id: Document ID

        Returns:
            List of reference dicts: [{"type": "citation", "text": "...", "position": 123}, ...]
        """
        references = []

        # Extract citation patterns (e.g., [1], [Smith 2020], (Smith et al., 2020))
        citation_patterns = [
            r"\[(\d+)\]",  # [1], [2], etc.
            r"\(([A-Z][a-z]+(?:\s+et\s+al\.)?,\s+\d{4})\)",  # (Smith, 2020) or (Smith et al., 2020)
            r"\[([A-Z][a-z]+(?:\s+et\s+al\.)?,\s+\d{4})\]",  # [Smith, 2020]
        ]

        for pattern in citation_patterns:
            for match in re.finditer(pattern, text):
                references.append(
                    {
                        "type": "citation",
                        "text": match.group(0),
                        "position": match.start(),
                        "document_id": document_id,
                    }
                )

        # Extract cross-references (e.g., "see Section 3", "refer to Chapter 2")
        cross_ref_patterns = [
            r"(?:see|refer to|see also)\s+(?:section|chapter|figure|table|appendix)\s+(\d+)",
            r"(?:see|refer to|see also)\s+(?:Section|Chapter|Figure|Table|Appendix)\s+(\d+)",
        ]

        for pattern in cross_ref_patterns:
            for match in re.finditer(pattern, text, re.IGNORECASE):
                references.append(
                    {
                        "type": "cross_reference",
                        "text": match.group(0),
                        "position": match.start(),
                        "document_id": document_id,
                    }
                )

        return references

    def extract_metadata(self, file_path: str, file_type: str) -> Dict[str, Any]:
        """
        Extract metadata from document.

        Args:
            file_path: Path to file
            file_type: MIME type

        Returns:
            Dict with metadata: title, author, creation_date, etc.
        """
        if file_type not in self.supported_types:
            return {}

        try:
            if file_type == "application/pdf":
                return self._extract_pdf_metadata(file_path)
            elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                return self._extract_docx_metadata(file_path)
            else:
                return {}
        except Exception as e:
            logger.warning(
                "[DocumentProcessor] Failed to extract metadata from %s: %s",
                file_path,
                e,
            )
            return {}

    def _extract_pdf_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from PDF."""
        metadata = {}
        if not PYPDF2_AVAILABLE:
            return metadata
        try:
            with open(file_path, "rb") as file:
                pdf_reader = PyPDF2.PdfReader(file)
                if pdf_reader.metadata:
                    pdf_meta = pdf_reader.metadata
                    if pdf_meta.get("/Title"):
                        metadata["title"] = pdf_meta["/Title"]
                    if pdf_meta.get("/Author"):
                        metadata["author"] = pdf_meta["/Author"]
                    if pdf_meta.get("/CreationDate"):
                        metadata["creation_date"] = str(pdf_meta["/CreationDate"])
                    if pdf_meta.get("/Subject"):
                        metadata["subject"] = pdf_meta["/Subject"]
        except Exception as e:
            logger.debug("[DocumentProcessor] PDF metadata extraction failed: %s", e)
        return metadata

    def _extract_docx_metadata(self, file_path: str) -> Dict[str, Any]:
        """Extract metadata from DOCX."""
        metadata = {}
        if not DOCX_AVAILABLE:
            return metadata
        try:
            doc = Document(file_path)
            core_props = doc.core_properties
            if core_props.title:
                metadata["title"] = core_props.title
            if core_props.author:
                metadata["author"] = core_props.author
            if core_props.created:
                metadata["creation_date"] = core_props.created.isoformat()
            if core_props.subject:
                metadata["subject"] = core_props.subject
            if core_props.keywords:
                metadata["keywords"] = core_props.keywords
        except Exception as e:
            logger.debug("[DocumentProcessor] DOCX metadata extraction failed: %s", e)
        return metadata

    def detect_language(self, text: str) -> Optional[str]:
        """
        Detect language of text.

        Args:
            text: Text to analyze

        Returns:
            Language code (e.g., 'zh', 'en', 'ja') or None if detection fails
        """
        # Ensure text is a string
        if isinstance(text, list):
            logger.warning("[DocumentProcessor] detect_language received list, converting to string")
            text = "\n".join(str(item) for item in text)
        if not isinstance(text, str):
            text = str(text) if text else ""

        if not LANGDETECT_AVAILABLE or not text or len(text.strip()) < 10:
            return None

        try:
            # Use first 1000 characters for faster detection
            sample = text[:1000] if len(text) > 1000 else text
            lang_code = detect(sample)
            logger.debug("[DocumentProcessor] Detected language: %s", lang_code)
            return lang_code
        except LangDetectException:
            return None
        except Exception as e:
            logger.debug("[DocumentProcessor] Language detection failed: %s", e)
            return None

    def extract_text(self, file_path: str, file_type: str) -> str:
        """
        Extract text from file.

        Args:
            file_path: Path to file
            file_type: MIME type

        Returns:
            Extracted text content
        """
        if file_type not in self.supported_types:
            raise ValueError(f"Unsupported file type: {file_type}")

        # Validate file content matches claimed type
        is_valid, detected_type = self.validate_file_content(file_path, file_type)
        if not is_valid:
            if detected_type:
                raise ValueError(
                    f"File content does not match claimed type. Claimed: {file_type}, Detected: {detected_type}"
                )
            else:
                raise ValueError(
                    f"File content validation failed. File may be corrupted or not match type: {file_type}"
                )

        extractor = self.supported_types[file_type]
        try:
            text = extractor(file_path)
            # Handle case where extractor might return a list (shouldn't happen, but defensive)
            if isinstance(text, list):
                logger.warning(
                    "[DocumentProcessor] Extractor returned list instead of string for %s, joining",
                    file_path,
                )
                text = "\n".join(str(item) for item in text)
            # Ensure text is a string
            if not isinstance(text, str):
                text = str(text) if text else ""
            if not text or not text.strip():
                raise ValueError(f"No text extracted from {file_path}")
            return text.strip()
        except Exception as e:
            logger.error("[DocumentProcessor] Failed to extract text from %s: %s", file_path, e)
            raise

    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF."""
        if not PYPDF2_AVAILABLE and not PDFPLUMBER_AVAILABLE:
            raise ImportError("PyPDF2 or pdfplumber required for PDF extraction")

        # Try PyPDF2 first if available
        if PYPDF2_AVAILABLE:
            try:
                text_parts = []
                with open(file_path, "rb") as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page in pdf_reader.pages:
                        text = page.extract_text()
                        if text:
                            text_parts.append(text)
                return "\n\n".join(text_parts)
            except Exception as e:
                logger.warning("[DocumentProcessor] PyPDF2 failed, trying pdfplumber: %s", e)
                # Fallback to pdfplumber if PyPDF2 fails
                if PDFPLUMBER_AVAILABLE:
                    with pdfplumber.open(file_path) as pdf:
                        text_parts = []
                        for page in pdf.pages:
                            text = page.extract_text()
                            if text:
                                text_parts.append(text)
                        return "\n\n".join(text_parts)
                raise

        # Use pdfplumber if PyPDF2 not available
        if PDFPLUMBER_AVAILABLE:
            with pdfplumber.open(file_path) as pdf:
                text_parts = []
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        text_parts.append(text)
                return "\n\n".join(text_parts)

        raise ImportError("PyPDF2 or pdfplumber required for PDF extraction")

    def extract_text_with_pages(self, file_path: str, file_type: str) -> tuple[str, List[Dict[str, Any]]]:
        """
        Extract text with page information for PDFs.

        Args:
            file_path: Path to file
            file_type: MIME type

        Returns:
            Tuple of (text, page_info) where page_info is list of page boundaries
        """
        if file_type == "application/pdf":
            if not PDFPLUMBER_AVAILABLE:
                # Fallback: extract text without page info
                text = self._extract_pdf(file_path)
                return text, []
            try:
                with pdfplumber.open(file_path) as pdf:
                    text_parts = []
                    page_info = []
                    current_pos = 0
                    for page_num, page in enumerate(pdf.pages, 1):
                        text = page.extract_text()
                        if text:
                            start_pos = current_pos
                            text_parts.append(text)
                            current_pos += len(text) + 2  # +2 for "\n\n"
                            page_info.append(
                                {
                                    "page": page_num,
                                    "start": start_pos,
                                    "end": current_pos,
                                }
                            )
                    return "\n\n".join(text_parts), page_info
            except ImportError:
                # Fallback: extract text without page info
                text = self._extract_pdf(file_path)
                return text, []
        else:
            # For non-PDF files, no page info
            text = self.extract_text(file_path, file_type)
            return text, []

    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX."""
        if not DOCX_AVAILABLE:
            raise ImportError("python-docx required for DOCX extraction")

        doc = Document(file_path)
        text_parts = []
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_parts.append(paragraph.text)
        return "\n\n".join(text_parts)

    def _extract_text(self, file_path: str) -> str:
        """Extract text from plain text files."""
        # Try multiple encodings
        encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as file:
                    return file.read()
            except UnicodeDecodeError:
                continue
            except Exception as e:
                logger.warning("[DocumentProcessor] Failed to read with %s: %s", encoding, e)
                continue

        raise ValueError(f"Failed to decode text file with any encoding: {file_path}")

    def _extract_image(self, file_path: str) -> str:
        """Extract text from image using OCR."""
        # Try DashScope OCR first
        try:
            return self._extract_image_dashscope(file_path)
        except Exception as e:
            logger.warning("[DocumentProcessor] DashScope OCR failed: %s", e)
            # Fallback to pytesseract
            try:
                return self._extract_image_tesseract(file_path)
            except Exception as e2:
                raise ValueError(f"OCR failed with both DashScope and Tesseract: {e2}") from e2

    def _extract_image_dashscope(self, file_path: str) -> str:
        """Extract text using DashScope OCR API."""
        if not HTTPX_AVAILABLE or not CONFIG_AVAILABLE:
            raise ValueError("httpx and config required for DashScope OCR")

        api_key = config.QWEN_API_KEY
        if not api_key:
            raise ValueError("DashScope API key required for OCR")

        # Read image and encode
        with open(file_path, "rb") as f:
            image_data = f.read()

        image_base64 = base64.b64encode(image_data).decode("utf-8")

        headers = {
            "Authorization": f"Bearer {api_key}",
            "Content-Type": "application/json",
        }

        payload = {
            "model": "qwen-vl-plus",
            "input": {
                "messages": [
                    {
                        "role": "user",
                        "content": [
                            {"image": f"data:image/jpeg;base64,{image_base64}"},
                            {"text": "请提取图片中的所有文字，保持原有格式。"},
                        ],
                    }
                ]
            },
            "parameters": {},
        }

        base_url = config.DASHSCOPE_API_URL or "https://dashscope.aliyuncs.com/api/v1/"
        ocr_url = f"{base_url}services/aigc/multimodal-generation/generation"

        with httpx.Client(timeout=60.0) as client:
            response = client.post(ocr_url, headers=headers, json=payload)
            response.raise_for_status()
            result = response.json()

            if "output" in result and "choices" in result["output"]:
                choices = result["output"]["choices"]
                if choices and len(choices) > 0:
                    content = choices[0].get("message", {}).get("content", "")
                    # Handle case where content might be a list (multimodal response)
                    if isinstance(content, list):
                        # Extract text from list of content blocks
                        text_parts = []
                        for item in content:
                            if isinstance(item, dict):
                                # Content block with type and text
                                if item.get("type") == "text":
                                    text_parts.append(item.get("text", ""))
                                elif "text" in item:
                                    text_parts.append(str(item["text"]))
                            elif isinstance(item, str):
                                text_parts.append(item)
                        text = "".join(text_parts)
                    elif isinstance(content, str):
                        text = content
                    else:
                        text = str(content) if content else ""

                    if text and text.strip():
                        return text.strip()

        raise ValueError("No text extracted from OCR response")

    def _extract_image_tesseract(self, file_path: str) -> str:
        """Extract text using pytesseract OCR."""
        if not TESSERACT_AVAILABLE:
            raise ImportError(
                "pytesseract and PIL required for Tesseract OCR. Install with: pip install pytesseract Pillow"
            )

        try:
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image, lang="chi_sim+eng")  # Chinese + English
            return text
        except Exception as e:
            # Check if this is a TesseractNotFoundError
            error_str = str(e).lower()
            if "tesseract" in error_str and (
                "not found" in error_str or "not installed" in error_str or "not in your path" in error_str
            ):
                logger.error("[DocumentProcessor] Tesseract OCR binary not found: %s", str(e))
                error_msg = (
                    "Tesseract OCR binary not found. "
                    "The pytesseract Python package is installed, but the Tesseract binary is missing or not in PATH.\n"
                    "Install Tesseract:\n"
                    "  - Ubuntu/Debian: sudo apt-get install tesseract-ocr tesseract-ocr-chi-sim\n"
                    "  - macOS: brew install tesseract\n"
                    "  - Windows: Download from https://github.com/UB-Mannheim/tesseract/wiki\n"
                    f"Original error: {str(e)}"
                )
                raise RuntimeError(error_msg) from e
            else:
                logger.error("[DocumentProcessor] Tesseract OCR failed: %s", str(e))
                error_msg = f"Tesseract OCR failed: {str(e)}"
                raise RuntimeError(error_msg) from e

    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX."""
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx required for PPTX extraction")

        prs = Presentation(file_path)
        text_parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text)
        return "\n\n".join(text_parts)

    def _extract_xlsx(self, file_path: str) -> str:
        """Extract text from XLSX."""
        if not OPENPYXL_AVAILABLE:
            raise ImportError("openpyxl required for XLSX extraction")

        wb = load_workbook(file_path, read_only=True)
        text_parts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = []
            for row in sheet.iter_rows(values_only=True):
                row_text = " ".join(str(cell) if cell else "" for cell in row)
                if row_text.strip():
                    sheet_text.append(row_text)
            if sheet_text:
                text_parts.append(f"Sheet: {sheet_name}\n" + "\n".join(sheet_text))
        return "\n\n".join(text_parts)

    def get_file_type(self, file_path: str) -> str:
        """
        Get MIME type of file.

        Args:
            file_path: Path to file

        Returns:
            MIME type string
        """
        mime_type, _ = mimetypes.guess_type(file_path)
        if not mime_type:
            # Fallback based on extension
            ext = Path(file_path).suffix.lower()
            ext_to_mime = {
                ".pdf": "application/pdf",
                ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                ".txt": "text/plain",
                ".md": "text/markdown",
                ".jpg": "image/jpeg",
                ".jpeg": "image/jpeg",
                ".png": "image/png",
                ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            }
            mime_type = ext_to_mime.get(ext, "application/octet-stream")

        return mime_type

    def is_supported(self, file_type: str) -> bool:
        """
        Check if file type is supported.

        Args:
            file_type: MIME type

        Returns:
            True if supported
        """
        return file_type in self.supported_types


def get_document_processor() -> DocumentProcessor:
    """Get global document processor instance."""
    if not hasattr(get_document_processor, "_instance"):
        get_document_processor._instance = DocumentProcessor()
    return get_document_processor._instance
